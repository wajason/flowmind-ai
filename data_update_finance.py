#!/usr/bin/env python3
"""
data_update_finance.py — FlowMind 知識庫入庫 pipeline
=============================================================================
沿用 AnalogGenie-RAG 已驗證的父子切塊 + HNSW/GIN 雙索引方法論，
針對供應鏈金融場域改了四處：

  1. 中文稀疏索引：fts_vector 改存字元 bigram（見 flowmind/textnorm.py）。
     原版對中文用 to_tsvector('english', …)，BM25 那一路實際上是壞的但不會報錯。

  2. Engagement 隔離：不再手寫 WHERE tenant_id，改由連線層 + PostgreSQL RLS 強制。
     這支程式即使寫錯 SQL，也寫不進別的客戶的資料。

  3. 多格式：新增 CSV / XLSX / DOCX / PPTX / JSON。
     統計表不逐列入庫（那是 CSV-in-RAG 的典型反模式，會讓檢索結果變成一堆
     破碎數字列），改用 pandas 產生結構化摘要；精確數字一律回原始檔案查。

  4. 分類標籤改為「法規 / 融資商品說明 / 白皮書統計 / 企業自有文件」，
     並可在檢索時按分類過濾 —— 因為「法規怎麼規定」和「這家客戶的發票長怎樣」
     在授信對話裡是兩種完全不同的問題。

Usage:
  python data_update_finance.py --tenant SHARED    --rebuild   # 建立公開知識庫
  python data_update_finance.py --tenant CASE-0001            # 匯入某客戶資料（增量）
  python data_update_finance.py --list                        # 列出所有委任案
"""

from __future__ import annotations

import argparse
import hashlib
import json
import re
import shutil
import sys
from pathlib import Path
from typing import Optional

sys.path.insert(0, str(Path(__file__).resolve().parent))

import psycopg2.extras

from flowmind import config, db, embeddings, textnorm

SUPPORTED_EXTS = {".pdf", ".md", ".txt", ".csv", ".xlsx", ".docx", ".pptx", ".json"}

# 依檔名關鍵字判斷分類。刻意用檔名而非內容判斷：規則透明、可預測，
# 使用者看得懂為什麼某份文件被分到某一類，也改得動。
# 規則由上往下比對，第一個命中者勝出，所以順序有意義：
# 「法規」放最前面，因為《商業會計法》這種檔名同時含「法」與「會計」，
# 若讓「融資商品說明」先比中就會被分錯類。
CATEGORY_RULES = [
    (["發展條例", "認定標準", "施行細則", "條例", "辦法", "準則", "法規",
      "民法", "營業稅法", "商業會計法", "保護法", "債權讓與", "銀行法",
      "票據法", "公司法"], "法規"),
    (["保證要點", "信用保證", "業務規章", "作業手冊", "承購", "factoring",
      "供應鏈金融", "融資", "貸款", "商品說明", "授信", "額度", "輔導"],
     "融資商品說明"),
    (["白皮書", "統計", "年報", "調查報告"], "白皮書統計"),
    (["發票", "invoice", "receivable", "payable", "contract", "合約", "ledger",
      "流水", "對帳", "customer_master", "cash_flow"], "企業自有文件"),
]


def guess_category(path: Path) -> str:
    name = path.stem.lower()
    for keywords, cat in CATEGORY_RULES:
        if any(k.lower() in name for k in keywords):
            return cat
    if path.suffix.lower() in (".csv", ".xlsx"):
        return "白皮書統計"
    return "企業自有文件"


# ── 資料來源登錄表（由 scripts/build_sources_manifest.py 產生）─────────────
_REGISTRY: dict[str, dict] | None = None


def load_registry() -> dict[str, dict]:
    """
    載入資料來源登錄表，取得每份文件的發布時間、狀態與權威層級。

    這是為了處理一個引用驗證擋不住的錯誤：**引用是真的，但答案過期了**。
    系統可以完全正確地引用 2015 年的舊作業手冊回答 2026 年的問題，
    而且引用驗證會顯示 exact 100 分 —— 因為那句話確實在那份文件裡。
    只能靠資料層的版本標註來擋。
    """
    global _REGISTRY
    if _REGISTRY is None:
        p = config.DATA_DIR / "sources_registry.json"
        if p.exists():
            data = json.loads(p.read_text(encoding="utf-8"))
            _REGISTRY = {e["filename"]: e for e in data.get("entries", [])}
        else:
            _REGISTRY = {}
            print("  [WARN] 找不到 data/sources_registry.json，"
                  "文件將缺少發布時間與版本狀態。"
                  "請先執行 python scripts/build_sources_manifest.py")
    return _REGISTRY


# ══════════════════════════════════════════════════════════════════════════
# 1. 文字抽取
# ══════════════════════════════════════════════════════════════════════════

def extract_pdf(path: Path) -> str:
    try:
        import pymupdf4llm
        return pymupdf4llm.to_markdown(str(path), show_progress=False)
    except ImportError:
        import fitz
        doc = fitz.open(str(path))
        return "\n\n".join(page.get_text() for page in doc)


def extract_md(path: Path) -> str:
    text = path.read_text(encoding="utf-8", errors="replace")
    return re.sub(r"!\[([^\]]*)\]\([^)]*\)", r"[圖: \1]", text)


def extract_txt(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


def extract_docx(path: Path) -> str:
    import docx
    doc = docx.Document(str(path))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n\n".join(parts)


def extract_pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        texts = [s.text_frame.text.strip() for s in slide.shapes
                 if s.has_text_frame and s.text_frame.text.strip()]
        if texts:
            parts.append(f"[投影片 {i}]\n" + "\n".join(texts))
    return "\n\n".join(parts)


def extract_csv_or_xlsx(path: Path) -> str:
    """
    統計表產生結構化摘要，不逐列入庫。
    RAG 需要知道的是「這份表在講什麼、量級多大」，
    精確數字要回原始檔案查 —— 讓 LLM 從破碎的數字列裡「記住」某個數值，
    正是幻覺最常發生的地方。
    """
    import pandas as pd
    try:
        df = (pd.read_csv(path, encoding="utf-8-sig") if path.suffix.lower() == ".csv"
              else pd.read_excel(path))
    except Exception as e:                            # noqa: BLE001
        return f"[無法解析統計表 {path.name}: {e}]"

    lines = [f"# 統計表摘要：{path.stem}", "",
             f"- 資料列數：{len(df)}　欄位數：{len(df.columns)}",
             f"- 欄位：{', '.join(str(c) for c in df.columns[:15])}"
             + (" …(更多欄位省略)" if len(df.columns) > 15 else "")]

    for col in df.select_dtypes(include="number").columns[:5]:
        s = df[col].dropna()
        if len(s):
            lines.append(f"- 「{col}」：範圍 {s.min():,.0f} ~ {s.max():,.0f}，"
                         f"總計 {s.sum():,.0f}，平均 {s.mean():,.0f}")
    # 不用 select_dtypes(include="object")：pandas 3 已把字串改成獨立的 str dtype，
    # 那個寫法會噴 deprecation warning 且未來會漏掉字串欄位。改用「非數值即類別」。
    numeric = set(df.select_dtypes(include="number").columns)
    for col in [c for c in df.columns if c not in numeric][:2]:
        top = df[col].value_counts().head(8)
        if len(top):
            lines.append(f"- 「{col}」主要類別：{', '.join(f'{k}({v})' for k, v in top.items())}")

    lines += ["", f"> 完整數據請查原始檔案 {path.name}。本摘要僅供語意檢索定位，"
                  f"不得作為精確數字的引用來源。"]
    return "\n".join(lines)


def extract_json(path: Path) -> str:
    """
    把結構化營運資料（發票/合約/現金流）轉成可檢索的自然語言。

    這一步不是為了讓 LLM 去「算」這些數字 —— 算數由 crosscheck.py 決定性完成。
    這裡的目的是讓「這家公司跟哪些買方往來、帳期多長、有沒有合約」
    這類語意層問題可以被檢索到，並且能引用回具體的憑證號碼。
    """
    data = json.loads(path.read_text(encoding="utf-8"))
    lines = [f"# {path.stem}", ""]

    if isinstance(data, dict):
        # 現金流預測這類單一物件：把時間軸壓成摘要，不逐筆展開
        if "timeline" in data:
            lines += [
                f"- 目前銀行餘額：NT${data.get('current_balance', 0):,.0f}",
                f"- 預測期間：未來 {data.get('horizon_days', 0)} 天",
                f"- 是否偵測到現金缺口：{'是' if data.get('gap_detected') else '否'}",
            ]
            if data.get("gap_detected"):
                lines.append(f"- 最早缺口日：{data['gap_date']}，"
                             f"缺口金額 NT${abs(data['gap_amount']):,.0f}")
            lines.append(f"- 逾期應收總額（不計入未來現金流）："
                         f"NT${data.get('overdue_receivables_total', 0):,.0f}")
            lines.append(f"- 計算方式：{data.get('computation_method')}（決定性運算，非 LLM 生成）")
            return "\n".join(lines)
        data = [data]

    if not isinstance(data, list) or not data:
        return json.dumps(data, ensure_ascii=False, indent=2)[:20000]

    doc_type = str(data[0].get("doc_type", "")) if isinstance(data[0], dict) else ""

    if doc_type == "AR_INVOICE":
        lines.append(f"本檔案為應收帳款發票明細，共 {len(data)} 張。")
        for inv in data:
            lines.append(
                f"- 發票 {inv['invoice_number']}（{inv['invoice_date']} 開立）："
                f"賣方 {inv['seller_name']}(統編 {inv['seller_ban']}) 開立予 "
                f"買方 {inv['buyer_name']}(統編 {inv['buyer_ban']})，"
                f"銷售額 NT${inv['sales_amount']:,}、稅額 NT${inv['tax_amount']:,}、"
                f"總計 NT${inv['total_amount']:,}，帳期 {inv['payment_terms_days']} 天，"
                f"到期日 {inv['due_date']}，狀態 {inv['status']}。")
    elif doc_type == "SALES_CONTRACT":
        lines.append(f"本檔案為年度基本買賣合約，共 {len(data)} 份。")
        for c in data:
            lines.append(
                f"- 合約 {c['contract_number']}：{c['seller_name']} 與 {c['buyer_name']}"
                f"(統編 {c['buyer_ban']})，有效期間 {c['effective_date']} 至 {c['expiry_date']}，"
                f"約定付款帳期 {c['payment_terms_days']} 天，"
                f"年度採購承諾金額 NT${c['annual_commitment_amount']:,}。{c.get('recourse_note','')}")
    elif doc_type == "AP_BILL":
        lines.append(f"本檔案為應付帳款明細，共 {len(data)} 筆。")
        for b in data:
            lines.append(
                f"- 應付 {b['bill_number']}（{b['issue_date']}）：應付予 {b['supplier_name']}"
                f"(統編 {b['supplier_ban']}) NT${b['amount']:,}，"
                f"帳期 {b['payment_terms_days']} 天，到期 {b['due_date']}，狀態 {b['status']}。")
    else:
        for item in data[:500]:
            lines.append("- " + "；".join(f"{k}={v}" for k, v in item.items()
                                          if k != "source_note"))
    return "\n".join(lines)


_EXTRACTORS = {
    ".pdf": extract_pdf, ".md": extract_md, ".txt": extract_txt,
    ".docx": extract_docx, ".pptx": extract_pptx,
    ".csv": extract_csv_or_xlsx, ".xlsx": extract_csv_or_xlsx,
    ".json": extract_json,
}


def extract_text(path: Path) -> Optional[str]:
    fn = _EXTRACTORS.get(path.suffix.lower())
    if fn is None:
        print(f"  [SKIP] 不支援的格式：{path.name}")
        return None
    try:
        return fn(path)
    except Exception as e:                            # noqa: BLE001
        print(f"  [ERROR] {path.name} 抽取失敗：{e}")
        return None


# ══════════════════════════════════════════════════════════════════════════
# 2. 清理與切塊（沿用 AnalogGenie 已驗證參數）
# ══════════════════════════════════════════════════════════════════════════

def clean_text(raw: str) -> str:
    protected: list[str] = []

    def keep(m):
        protected.append(m.group(0))
        return f"__PROTECTED_{len(protected)-1}__"

    text = re.sub(r"```.*?```|\$\$.*?\$\$", keep, raw, flags=re.DOTALL)
    text = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]", "", text)
    text = re.sub(r"(?m)^[-–]\s*\d+\s*[-–]\s*$", "", text)
    text = re.sub(r"(?im)^page\s+\d+\s+(of\s+\d+)?\s*$", "", text)
    text = re.sub(r"(?m)^第\s*\d+\s*頁(，共\s*\d+\s*頁)?\s*$", "", text)   # 中文頁碼
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r"(?m)^[=\-_]{5,}\s*$", "", text)
    text = "\n".join(line.rstrip() for line in text.split("\n"))
    for i, block in enumerate(protected):
        text = text.replace(f"__PROTECTED_{i}__", block)
    return text.strip()


def chunk_text(text: str, source: str, category: str,
               registry_entry: dict | None = None) -> list[dict]:
    from langchain_text_splitters import (MarkdownHeaderTextSplitter,
                                          RecursiveCharacterTextSplitter)

    md_splitter = MarkdownHeaderTextSplitter(
        headers_to_split_on=[("#", "H1"), ("##", "H2"), ("###", "H3")])
    md_docs = md_splitter.split_text(text)

    parent_fallback = RecursiveCharacterTextSplitter(
        chunk_size=config.PARENT_FALLBACK_SIZE, chunk_overlap=config.PARENT_OVERLAP,
        separators=["\n\n##", "\n\n", "\n", "。", " "])
    child_splitter = RecursiveCharacterTextSplitter(
        chunk_size=config.CHILD_CHUNK_SIZE, chunk_overlap=config.CHILD_CHUNK_OVERLAP,
        # 中文用「。」「；」當分隔點才切得出完整句子；只用 "." 會切在小數點中間
        separators=["\n\n", "\n", "。", "；", ".", " "])

    chunks, idx = [], 0
    for md_doc in md_docs:
        headers = [v for k, v in md_doc.metadata.items() if k.startswith("H")]
        header_ctx = " > ".join(headers) if headers else "全文脈絡"

        if len(md_doc.page_content) <= config.PARENT_MAX_SIZE:
            parents, intact = [md_doc], True
        else:
            parents, intact = parent_fallback.split_documents([md_doc]), False

        for p in parents:
            parent_text = f"[分類: {category}] [章節: {header_ctx}]\n{p.page_content}"
            for c in child_splitter.split_documents([p]):
                meta = {"category": category, "headers": headers,
                        "parent_content": parent_text,
                        "is_intact_section": intact}
                if registry_entry:
                    meta.update({
                        "published": registry_entry.get("published"),
                        "doc_status": registry_entry.get("status"),
                        "authority": registry_entry.get("authority"),
                        "publisher": registry_entry.get("publisher"),
                        "superseded_by": registry_entry.get("superseded_by"),
                    })
                chunks.append({
                    "content": f"[{header_ctx}]\n{c.page_content}",
                    "source": source, "chunk_index": idx, "metadata": meta,
                })
                idx += 1
    return chunks


# ══════════════════════════════════════════════════════════════════════════
# 3. 入庫
# ══════════════════════════════════════════════════════════════════════════

def delete_by_source(conn, source: str) -> None:
    # 沒有 WHERE tenant_id —— RLS 會自動限縮在目前 engagement 內。
    # 這正是把隔離下沉到資料庫的價值：連「刪錯別人資料」都不可能發生。
    with conn.cursor() as cur:
        cur.execute("DELETE FROM documents WHERE source = %s", (source,))
    conn.commit()


def insert_chunks(conn, tenant_id: str, chunks: list[dict], file_hash: str) -> None:
    vectors = embeddings.embed([c["content"] for c in chunks])
    rows = []
    for c, v in zip(chunks, vectors):
        meta = c["metadata"]
        # fts 索引的來源是 parent_content：child 只有 400 字元，
        # 關鍵詞常常落在 child 之外、卻在同一個 parent 裡。
        fts_src = f"{meta.get('category','')} {meta.get('parent_content','')}"
        rows.append((tenant_id, c["source"], c["chunk_index"], c["content"],
                     embeddings.to_pgvector(v), json.dumps(meta, ensure_ascii=False),
                     textnorm.to_fts_document(fts_src), file_hash))

    with conn.cursor() as cur:
        psycopg2.extras.execute_values(cur, """
            INSERT INTO documents
              (tenant_id, source, chunk_index, content, embedding, metadata, fts_vector, file_hash)
            VALUES %s
            ON CONFLICT (tenant_id, source, chunk_index) DO UPDATE SET
              content=EXCLUDED.content, embedding=EXCLUDED.embedding,
              metadata=EXCLUDED.metadata, fts_vector=EXCLUDED.fts_vector,
              file_hash=EXCLUDED.file_hash, created_at=NOW();
        """, rows, template="(%s,%s,%s,%s,%s::vector,%s::jsonb,to_tsvector('simple',%s),%s)")
    conn.commit()


# ══════════════════════════════════════════════════════════════════════════
# 4. 檔案指紋追蹤
# ══════════════════════════════════════════════════════════════════════════

def hash_path(tenant_id: str) -> Path:
    return config.DATA_DIR / ".file_hashes" / f"{tenant_id}.json"


def load_hashes(tenant_id: str) -> dict[str, str]:
    p = hash_path(tenant_id)
    return json.loads(p.read_text(encoding="utf-8")) if p.exists() else {}


def save_hashes(tenant_id: str, hashes: dict[str, str]) -> None:
    p = hash_path(tenant_id)
    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_text(json.dumps(hashes, indent=2, ensure_ascii=False), encoding="utf-8")


def file_hash(path: Path) -> str:
    h = hashlib.sha256()
    with path.open("rb") as f:
        for block in iter(lambda: f.read(1 << 20), b""):
            h.update(block)
    return h.hexdigest()


# ══════════════════════════════════════════════════════════════════════════
# 5. 主流程
# ══════════════════════════════════════════════════════════════════════════

def process_file(path: Path, conn, tenant_id: str, processed_dir: Path) -> bool:
    raw = extract_text(path)
    if raw is None:
        return False
    cleaned = clean_text(raw)
    if len(cleaned) < 50:
        print(f"    [WARN] 清理後內容過短，跳過")
        return False

    processed_dir.mkdir(parents=True, exist_ok=True)
    (processed_dir / (path.stem + ".txt")).write_text(cleaned, encoding="utf-8")

    entry = load_registry().get(path.name)
    # 登錄表的分類優先於檔名猜測 —— 登錄表是人工確認過的，檔名規則只是 fallback
    category = (entry or {}).get("category") or guess_category(path)
    chunks = chunk_text(cleaned, source=path.name, category=category,
                        registry_entry=entry)
    status = (entry or {}).get("status", "未登錄")
    pub = (entry or {}).get("published") or "—"
    print(f"    → {len(chunks)} chunks｜分類：{category}｜發布：{pub}｜狀態：{status}")

    delete_by_source(conn, path.name)
    insert_chunks(conn, tenant_id, chunks, file_hash(path))
    return True


def run(tenant_id: str, rebuild: bool) -> None:
    raw_dir = config.RAW_DIR / tenant_id
    processed_dir = config.PROCESSED_DIR / tenant_id

    if not raw_dir.exists():
        print(f"錯誤：{raw_dir} 不存在。請建立此資料夾並放入該 engagement 的文件。")
        sys.exit(1)

    # engagement 必須先登記才能入庫（外鍵約束）。這不是官僚流程 ——
    # 它強迫每一批資料都有明確的委任來源與保存期限，這是金融場域的基本要求。
    db.upsert_engagement(tenant_id, tenant_id, "未指定" if tenant_id != db.SHARED
                         else "法規與融資商品公開資料")

    with db.tenant_session(tenant_id) as conn:
        if rebuild:
            print(f"\n🔄 重建模式（engagement={tenant_id}）\n")
            if processed_dir.exists():
                shutil.rmtree(processed_dir)
            with conn.cursor() as cur:
                cur.execute("DELETE FROM documents")   # RLS 限縮在本 engagement
            conn.commit()
            stored = {}
        else:
            print(f"\n⚡ 增量更新（engagement={tenant_id}）\n")
            stored = load_hashes(tenant_id)

        current = {f.name: f for f in raw_dir.iterdir()
                   if f.suffix.lower() in SUPPORTED_EXTS and not f.name.startswith("_")}

        for name in set(stored) - set(current):
            print(f"  🗑️  移除已刪除的檔案：{name}")
            delete_by_source(conn, name)
            del stored[name]

        changed = skipped = 0
        for name, path in sorted(current.items()):
            h = file_hash(path)
            if stored.get(name) == h:
                skipped += 1
                continue
            print(f"  📄 {'新增' if name not in stored else '異動'}：{name}")
            if process_file(path, conn, tenant_id, processed_dir):
                stored[name] = h
                changed += 1

        save_hashes(tenant_id, stored)

        with conn.cursor() as cur:
            cur.execute("SELECT COUNT(*), COUNT(DISTINCT source) FROM documents")
            n_chunks, n_docs = cur.fetchone()
        db.write_audit(conn, tenant_id=tenant_id, action="ingest",
                       query_text=f"rebuild={rebuild}",
                       doc_sources=sorted(current.keys())[:50])

    print(f"\n✅ 完成 [{tenant_id}]：更新 {changed}、未變動 {skipped}")
    print(f"   目前此 engagement 可檢索：{n_docs} 份文件 / {n_chunks} 個 chunk")


def main():
    ap = argparse.ArgumentParser(description="FlowMind 知識庫入庫 pipeline")
    ap.add_argument("--tenant", help="'SHARED' 為公開知識庫；'CASE-xxxx' 為個別委任案")
    ap.add_argument("--rebuild", action="store_true")
    ap.add_argument("--list", action="store_true", help="列出所有 engagement 與其文件數")
    args = ap.parse_args()

    if args.list:
        print(f"\n{'engagement':<14}{'客戶':<28}{'類型':<24}{'文件':>6}{'chunks':>9}")
        print("─" * 84)
        for e in db.list_engagements():
            print(f"{e['tenant_id']:<14}{(e['client_name'] or '')[:26]:<28}"
                  f"{(e['engagement_type'] or '')[:22]:<24}{e['docs']:>6}{e['chunks']:>9}")
        return

    if not args.tenant:
        ap.error("需要 --tenant 或 --list")
    run(args.tenant, args.rebuild)


if __name__ == "__main__":
    main()
